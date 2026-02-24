"""
Recast - File Converter
A free, local file converter for Windows.
Supports images, documents, audio, video, ebooks and archives.
"""

import customtkinter as ctk
from tkinter import filedialog, messagebox
import os
import threading
import subprocess
import zipfile
import tarfile
import json
import xml.etree.ElementTree as ET
import re

from PIL import Image

try:
    import pillow_heif
    pillow_heif.register_heif_opener()
    HEIC_AVAILABLE = True
except ImportError:
    HEIC_AVAILABLE = False

import pandas as pd
from docx import Document
from fpdf import FPDF
from pptx import Presentation
import markdown
from pydub import AudioSegment
import ffmpeg

try:
    import py7zr
    SEVENZIP_AVAILABLE = True
except ImportError:
    SEVENZIP_AVAILABLE = False


# ---------------------------------------------------------------------------
# Optional tool detection (runs once at startup)
# ---------------------------------------------------------------------------

def _check_tool(cmd):
    try:
        subprocess.run([cmd, "--version"], capture_output=True, timeout=5)
        return True
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False

LIBREOFFICE_AVAILABLE = _check_tool("libreoffice") or _check_tool("soffice")
CALIBRE_AVAILABLE     = _check_tool("ebook-convert")

# Conversions that require LibreOffice: {in_fmt: [out_fmts]}
REQUIRES_LIBREOFFICE = {
    "pdf":  ["docx", "txt", "html", "rtf"],
    "docx": ["pdf", "rtf", "odt"],
    "html": ["pdf", "docx", "rtf"],
    "pptx": ["pdf"],
    "odt":  ["pdf", "docx", "txt", "html", "rtf"],
    "rtf":  ["pdf", "docx", "txt", "html"],
}

LIBREOFFICE_INSTALL = "Requires LibreOffice\nlibreoffice.org — free install"
CALIBRE_INSTALL     = "Requires Calibre\ncalibre-ebook.com — free install"


def get_chip_disabled_reason(category, in_fmt, out_fmt):
    """Returns tooltip string if conversion unavailable, else None."""
    if category == "ebook" and not CALIBRE_AVAILABLE:
        return CALIBRE_INSTALL
    if out_fmt in REQUIRES_LIBREOFFICE.get(in_fmt, []) and not LIBREOFFICE_AVAILABLE:
        return LIBREOFFICE_INSTALL
    return None


# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

BG         = "#0d0d0f"
SIDEBAR_BG = "#111114"
CARD_BG    = "#18181c"
CARD_HOVER = "#1f1f25"
BORDER     = "#2a2a32"
ACCENT     = "#6c63ff"
ACCENT2    = "#a78bfa"
SUCCESS    = "#22d3a0"
WARNING    = "#f59e0b"
ERROR      = "#f87171"
TEXT_PRI   = "#f0f0f5"
TEXT_SEC   = "#a0a0c0"
TEXT_DIM   = "#6a6a88"

FONT_HEAD  = ("Segoe UI", 22, "bold")
FONT_TITLE = ("Segoe UI", 13, "bold")
FONT_BODY  = ("Segoe UI", 12)
FONT_SMALL = ("Segoe UI", 10)
FONT_MONO  = ("Consolas", 11)

ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")


# ---------------------------------------------------------------------------
# Format data
# ---------------------------------------------------------------------------

CATEGORIES = [
    {"id": "all",      "icon": "⬡",  "label": "All Files"},
    {"id": "image",    "icon": "🖼",  "label": "Images"},
    {"id": "document", "icon": "📄",  "label": "Documents"},
    {"id": "audio",    "icon": "🎵",  "label": "Audio"},
    {"id": "video",    "icon": "🎬",  "label": "Video"},
    {"id": "ebook",    "icon": "📚",  "label": "Ebooks"},
    {"id": "archive",  "icon": "📦",  "label": "Archives"},
]

FORMAT_CATEGORIES = {
    "image":    ["png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff", "ico",
                 "heic", "heif", "tga", "ppm"],
    "document": ["pdf", "docx", "txt", "csv", "xlsx", "xls", "pptx",
                 "rtf", "html", "md", "json", "xml", "odt"],
    "audio":    ["mp3", "wav", "flac", "aac", "ogg", "m4a", "wma", "opus", "aiff"],
    "video":    ["mp4", "mov", "avi", "mkv", "webm", "flv", "ts",
                 "m4v", "3gp", "wmv", "mpeg", "mpg"],
    "ebook":    ["epub", "mobi", "azw3", "fb2"],
    "archive":  ["zip", "tar", "gz", "7z", "bz2"],
}

OUTPUT_FORMATS = {
    "image":    ["png", "jpg", "webp", "gif", "bmp", "tiff", "ico", "tga", "ppm"],
    "document": ["pdf", "docx", "txt", "csv", "xlsx", "html", "md", "rtf", "json", "xml"],
    "audio":    ["mp3", "wav", "flac", "ogg", "m4a", "aac", "opus", "aiff"],
    "video":    ["mp4", "avi", "mkv", "webm", "mov", "gif", "mp3"],
    "ebook":    ["epub", "mobi", "azw3", "txt", "pdf"],
    "archive":  ["zip", "tar", "gz", "7z"],
}

FORMAT_GROUPS = {
    "image": {
        "🖼 Raster":  ["png", "jpg", "webp", "gif", "bmp", "tiff"],
        "🔷 Special": ["ico", "tga", "ppm"],
    },
    "document": {
        "📝 Text":   ["txt", "md", "rtf", "html"],
        "📊 Data":   ["csv", "xlsx", "json", "xml"],
        "📄 Office": ["pdf", "docx"],
    },
    "audio": {
        "🎵 Lossy":    ["mp3", "aac", "ogg", "m4a", "opus"],
        "🎼 Lossless": ["wav", "flac", "aiff"],
    },
    "video": {
        "🎬 Video":   ["mp4", "avi", "mkv", "webm", "mov"],
        "✨ Special": ["gif", "mp3"],
    },
    "ebook":   {"📚 Ebook":   ["epub", "mobi", "azw3", "txt", "pdf"]},
    "archive": {"📦 Archive": ["zip", "tar", "gz", "7z"]},
}


def get_category(ext):
    ext = ext.lower().lstrip(".")
    for cat, exts in FORMAT_CATEGORIES.items():
        if ext in exts:
            return cat
    return None


def format_filesize(bytes_):
    for unit in ["B", "KB", "MB", "GB"]:
        if bytes_ < 1024:
            return f"{bytes_:.1f} {unit}"
        bytes_ /= 1024
    return f"{bytes_:.1f} TB"


# ---------------------------------------------------------------------------
# Converters
# ---------------------------------------------------------------------------

def convert_image(src, dst, out_fmt):
    img = Image.open(src)
    fmt_map = {
        "jpg": "JPEG", "jpeg": "JPEG", "tiff": "TIFF",
        "tga": "TGA",  "ico":  "ICO",  "ppm":  "PPM",
    }
    save_fmt = fmt_map.get(out_fmt.lower(), out_fmt.upper())

    if save_fmt == "ICO":
        img = img.resize((256, 256), Image.LANCZOS)

    if save_fmt in ("JPEG", "BMP", "PPM") and img.mode in ("RGBA", "P", "LA"):
        bg = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode == "P":
            img = img.convert("RGBA")
        bg.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
        img = bg
    elif save_fmt not in ("GIF", "PNG", "WEBP", "TIFF") and img.mode == "P":
        img = img.convert("RGB")

    img.save(dst, format=save_fmt)


def convert_audio(src, dst, out_fmt):
    try:
        audio = AudioSegment.from_file(src)
        codec_map = {"ogg": "libvorbis", "opus": "libopus", "aac": "aac", "m4a": "aac"}
        kwargs = {"format": out_fmt}
        if out_fmt in codec_map:
            kwargs["codec"] = codec_map[out_fmt]
        audio.export(dst, **kwargs)
    except Exception:
        ffmpeg.input(src).output(dst).overwrite_output().run(quiet=True)


def convert_video(src, dst, out_fmt):
    if out_fmt == "gif":
        palette = dst.replace(".gif", "_palette.png")
        try:
            (ffmpeg.input(src)
             .filter("fps", fps=10)
             .filter("scale", 480, -1, flags="lanczos")
             .filter("palettegen")
             .output(palette, vframes=1)
             .overwrite_output().run(quiet=True))
            in1 = ffmpeg.input(src).filter("fps", fps=10).filter("scale", 480, -1, flags="lanczos")
            in2 = ffmpeg.input(palette)
            (ffmpeg.filter([in1, in2], "paletteuse")
             .output(dst).overwrite_output().run(quiet=True))
        finally:
            if os.path.exists(palette):
                os.remove(palette)
    elif out_fmt == "mp3":
        (ffmpeg.input(src)
         .output(dst, vn=None, acodec="libmp3lame")
         .overwrite_output().run(quiet=True))
    else:
        ffmpeg.input(src).output(dst).overwrite_output().run(quiet=True)


def convert_document(src, dst, in_fmt, out_fmt):
    if in_fmt in ("csv", "xlsx", "xls"):
        df = pd.read_csv(src) if in_fmt == "csv" else pd.read_excel(src)
        handlers = {
            "xlsx": lambda: df.to_excel(dst, index=False),
            "csv":  lambda: df.to_csv(dst, index=False),
            "json": lambda: df.to_json(dst, orient="records", indent=2),
            "html": lambda: df.to_html(dst, index=False),
            "xml":  lambda: df.to_xml(dst, index=False),
            "txt":  lambda: open(dst, "w", encoding="utf-8").write(df.to_string()),
            "pdf":  lambda: _txt_to_pdf(df.to_string(), dst),
        }
        if out_fmt not in handlers:
            raise NotImplementedError(f"{in_fmt} → {out_fmt} not supported.")
        handlers[out_fmt]()
        return

    if in_fmt == "txt":
        text = open(src, "r", encoding="utf-8", errors="ignore").read()
        if out_fmt == "pdf":       _txt_to_pdf(text, dst)
        elif out_fmt == "docx":    _txt_to_docx(text, dst)
        elif out_fmt == "rtf":     _txt_to_rtf(text, dst)
        elif out_fmt == "html":
            open(dst, "w", encoding="utf-8").write(f"<html><body><pre>{text}</pre></body></html>")
        elif out_fmt == "md":
            open(dst, "w", encoding="utf-8").write(text)
        elif out_fmt == "xlsx":
            pd.DataFrame({"Content": text.splitlines()}).to_excel(dst, index=False)
        elif out_fmt == "csv":
            pd.DataFrame({"Content": text.splitlines()}).to_csv(dst, index=False)
        else:
            raise NotImplementedError(f"TXT → {out_fmt} not supported.")
        return

    if in_fmt == "md":
        text = open(src, "r", encoding="utf-8", errors="ignore").read()
        if out_fmt == "html":
            html = markdown.markdown(text, extensions=["tables", "fenced_code"])
            open(dst, "w", encoding="utf-8").write(
                f"<html><head><meta charset='utf-8'></head><body>{html}</body></html>"
            )
        elif out_fmt == "txt":
            open(dst, "w", encoding="utf-8").write(re.sub(r"[#*`>\[\]_~]", "", text))
        elif out_fmt == "pdf":     _txt_to_pdf(text, dst)
        elif out_fmt == "docx":    _txt_to_docx(text, dst)
        else:
            raise NotImplementedError(f"MD → {out_fmt} not supported.")
        return

    if in_fmt == "docx":
        doc = Document(src)
        if out_fmt == "txt":
            with open(dst, "w", encoding="utf-8") as f:
                for p in doc.paragraphs:
                    f.write(p.text + "\n")
        elif out_fmt == "html":
            with open(dst, "w", encoding="utf-8") as f:
                f.write("<html><body>" + "".join(f"<p>{p.text}</p>" for p in doc.paragraphs) + "</body></html>")
        elif out_fmt == "md":
            with open(dst, "w", encoding="utf-8") as f:
                for p in doc.paragraphs:
                    s = p.style.name.lower()
                    prefix = "# " if "heading 1" in s else "## " if "heading 2" in s else "### " if "heading 3" in s else ""
                    f.write(f"{prefix}{p.text}\n\n")
        elif out_fmt in ("pdf", "rtf", "odt"):
            _libreoffice_convert(src, dst, out_fmt)
        else:
            raise NotImplementedError(f"DOCX → {out_fmt} not supported.")
        return

    if in_fmt == "json":
        data = json.load(open(src, "r", encoding="utf-8"))
        if out_fmt in ("csv", "xlsx"):
            if not isinstance(data, list):
                raise ValueError("JSON must be an array of objects to convert to CSV/XLSX.")
            df = pd.DataFrame(data)
            df.to_csv(dst, index=False) if out_fmt == "csv" else df.to_excel(dst, index=False)
        elif out_fmt == "txt":
            open(dst, "w", encoding="utf-8").write(json.dumps(data, indent=2))
        elif out_fmt == "xml":
            root = ET.Element("root")
            for item in (data if isinstance(data, list) else [data]):
                child = ET.SubElement(root, "item")
                if isinstance(item, dict):
                    for k, v in item.items():
                        ET.SubElement(child, str(k)).text = str(v)
            ET.ElementTree(root).write(dst, encoding="unicode", xml_declaration=True)
        else:
            raise NotImplementedError(f"JSON → {out_fmt} not supported.")
        return

    if in_fmt == "xml":
        tree = ET.parse(src)
        if out_fmt == "txt":
            ET.indent(tree)
            tree.write(dst, encoding="unicode")
        elif out_fmt == "json":
            def el_to_dict(el):
                result = {}
                for child in el:
                    val = el_to_dict(child) if len(child) else child.text
                    result.setdefault(child.tag, []).append(val)
                return {k: v[0] if len(v) == 1 else v for k, v in result.items()} or el.text
            open(dst, "w", encoding="utf-8").write(
                json.dumps(el_to_dict(tree.getroot()), indent=2)
            )
        else:
            raise NotImplementedError(f"XML → {out_fmt} not supported.")
        return

    if in_fmt == "html":
        text = open(src, "r", encoding="utf-8", errors="ignore").read()
        if out_fmt == "txt":
            open(dst, "w", encoding="utf-8").write(re.sub(r"<[^>]+>", "", text))
        elif out_fmt == "md":
            t = re.sub(r"<h1[^>]*>(.*?)</h1>", r"# \1\n", text, flags=re.DOTALL)
            t = re.sub(r"<h2[^>]*>(.*?)</h2>", r"## \1\n", t,    flags=re.DOTALL)
            t = re.sub(r"<p[^>]*>(.*?)</p>",   r"\1\n\n", t,     flags=re.DOTALL)
            t = re.sub(r"<strong>(.*?)</strong>", r"**\1**", t)
            t = re.sub(r"<em>(.*?)</em>",         r"*\1*",  t)
            t = re.sub(r"<[^>]+>", "", t)
            open(dst, "w", encoding="utf-8").write(t)
        elif out_fmt in ("pdf", "docx", "rtf"):
            _libreoffice_convert(src, dst, out_fmt)
        else:
            raise NotImplementedError(f"HTML → {out_fmt} not supported.")
        return

    if in_fmt == "pdf":
        if out_fmt in ("docx", "txt", "html", "rtf"):
            _libreoffice_convert(src, dst, out_fmt)
        else:
            raise NotImplementedError(f"PDF → {out_fmt} not supported.")
        return

    if in_fmt == "pptx":
        prs = Presentation(src)
        if out_fmt == "txt":
            with open(dst, "w", encoding="utf-8") as f:
                for i, slide in enumerate(prs.slides, 1):
                    f.write(f"\n--- Slide {i} ---\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            f.write(shape.text + "\n")
        elif out_fmt == "html":
            slides_html = ""
            for i, slide in enumerate(prs.slides, 1):
                content = "".join(
                    f"<p>{s.text}</p>" for s in slide.shapes if hasattr(s, "text")
                )
                slides_html += f"<div><h2>Slide {i}</h2>{content}</div><hr>"
            open(dst, "w", encoding="utf-8").write(f"<html><body>{slides_html}</body></html>")
        elif out_fmt == "pdf":
            _libreoffice_convert(src, dst, "pdf")
        else:
            raise NotImplementedError(f"PPTX → {out_fmt} not supported.")
        return

    if in_fmt in ("odt", "rtf"):
        _libreoffice_convert(src, dst, out_fmt)
        return

    raise NotImplementedError(f"{in_fmt.upper()} → {out_fmt.upper()} not supported.")


def convert_ebook(src, dst):
    try:
        subprocess.run(["ebook-convert", src, dst], check=True, capture_output=True)
    except FileNotFoundError:
        raise RuntimeError("Ebook conversion requires Calibre.\nhttps://calibre-ebook.com")


def convert_archive(src, dst, in_fmt, out_fmt):
    import tempfile
    import shutil

    tmp = tempfile.mkdtemp()
    try:
        if in_fmt == "zip":
            with zipfile.ZipFile(src, "r") as z:
                z.extractall(tmp)
        elif in_fmt in ("tar", "gz", "bz2"):
            with tarfile.open(src) as t:
                t.extractall(tmp)
        elif in_fmt == "7z":
            if not SEVENZIP_AVAILABLE:
                raise RuntimeError("7z support requires: pip install py7zr")
            with py7zr.SevenZipFile(src, mode="r") as z:
                z.extractall(tmp)

        if out_fmt == "zip":
            with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as z:
                for root, _, files in os.walk(tmp):
                    for file in files:
                        fp = os.path.join(root, file)
                        z.write(fp, os.path.relpath(fp, tmp))
        elif out_fmt in ("tar", "gz", "bz2"):
            mode = {"tar": "w", "gz": "w:gz", "bz2": "w:bz2"}[out_fmt]
            with tarfile.open(dst, mode) as t:
                t.add(tmp, arcname="")
        elif out_fmt == "7z":
            if not SEVENZIP_AVAILABLE:
                raise RuntimeError("7z support requires: pip install py7zr")
            with py7zr.SevenZipFile(dst, mode="w") as z:
                z.writeall(tmp)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ---------------------------------------------------------------------------
# Document helpers
# ---------------------------------------------------------------------------

def _txt_to_pdf(text, dst):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    for line in text.splitlines():
        pdf.multi_cell(0, 6, line)
    pdf.output(dst)


def _txt_to_docx(text, dst):
    doc = Document()
    for line in text.splitlines():
        doc.add_paragraph(line)
    doc.save(dst)


def _txt_to_rtf(text, dst):
    escaped = text.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}")
    rtf_lines = "\n".join(f"\\pard {line}\\par" for line in escaped.splitlines())
    with open(dst, "w", encoding="ascii", errors="ignore") as f:
        f.write("{\\rtf1\\ansi\n" + rtf_lines + "\n}")


def _libreoffice_convert(src, dst, out_fmt):
    out_dir = os.path.dirname(os.path.abspath(dst))
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", out_fmt, "--outdir", out_dir, src],
            check=True, capture_output=True, text=True,
        )
        base   = os.path.splitext(os.path.basename(src))[0]
        lo_out = os.path.join(out_dir, f"{base}.{out_fmt}")
        if lo_out != dst and os.path.exists(lo_out):
            os.replace(lo_out, dst)
    except FileNotFoundError:
        raise RuntimeError("This conversion requires LibreOffice.\nhttps://www.libreoffice.org")


# ---------------------------------------------------------------------------
# UI components
# ---------------------------------------------------------------------------

class SidebarButton(ctk.CTkFrame):
    def __init__(self, master, icon, label, cat_id, on_click, **kwargs):
        super().__init__(master, fg_color="transparent", cursor="hand2", **kwargs)
        self.cat_id   = cat_id
        self.on_click = on_click
        self.active   = False

        self.inner = ctk.CTkFrame(self, fg_color="transparent", corner_radius=10)
        self.inner.pack(fill="x", padx=8, pady=2)

        self.icon_lbl = ctk.CTkLabel(
            self.inner, text=icon,
            font=("Segoe UI Emoji", 16), text_color=TEXT_SEC, width=28,
        )
        self.icon_lbl.pack(side="left", padx=(12, 6), pady=10)

        self.text_lbl = ctk.CTkLabel(
            self.inner, text=label,
            font=FONT_BODY, text_color=TEXT_SEC, anchor="w",
        )
        self.text_lbl.pack(side="left", fill="x", expand=True)

        for w in [self, self.inner, self.icon_lbl, self.text_lbl]:
            w.bind("<Button-1>", self._click)
            w.bind("<Enter>",    self._hover_in)
            w.bind("<Leave>",    self._hover_out)

    def _click(self, e=None):
        self.on_click(self.cat_id)

    def _hover_in(self, e=None):
        if not self.active:
            self.inner.configure(fg_color=CARD_HOVER)

    def _hover_out(self, e=None):
        if not self.active:
            self.inner.configure(fg_color="transparent")

    def set_active(self, active):
        self.active = active
        color = ACCENT2 if active else TEXT_SEC
        bg    = CARD_BG if active else "transparent"
        self.inner.configure(fg_color=bg)
        self.icon_lbl.configure(text_color=color)
        self.text_lbl.configure(text_color=color)


class FormatChip(ctk.CTkFrame):
    def __init__(self, master, fmt, on_select, disabled_reason=None, **kwargs):
        bg = "#111116" if disabled_reason else CARD_BG
        super().__init__(master, fg_color=bg, corner_radius=10,
                         border_width=1,
                         border_color=BORDER if disabled_reason else "#2e2e3a",
                         cursor="hand2" if not disabled_reason else "arrow", **kwargs)
        self.fmt             = fmt
        self.on_select       = on_select
        self.selected        = False
        self.disabled_reason = disabled_reason
        self._tooltip_win    = None

        inner = ctk.CTkFrame(self, fg_color="transparent")
        inner.pack(padx=14, pady=10)

        # Format name — large and readable
        name_color = "#404058" if disabled_reason else "#d0d0e8"
        self.lbl = ctk.CTkLabel(
            inner, text=fmt.upper(),
            font=("Segoe UI", 13, "bold"),
            text_color=name_color,
        )
        self.lbl.pack()

        # Extension tag below — clearly readable
        if not disabled_reason:
            self.sub_lbl = ctk.CTkLabel(
                inner, text=f".{fmt}",
                font=("Consolas", 12, "bold"),
                text_color="#9090b8",
            )
        else:
            self.sub_lbl = ctk.CTkLabel(
                inner, text="🔒 install required",
                font=("Segoe UI", 9),
                text_color="#404058",
            )
        self.sub_lbl.pack(pady=(1, 0))

        for w in [self, inner, self.lbl, self.sub_lbl]:
            w.bind("<Button-1>", self._click)
            w.bind("<Enter>",    self._hover_in)
            w.bind("<Leave>",    self._hover_out)

    def _click(self, e=None):
        if not self.disabled_reason:
            self.on_select(self.fmt)

    def _hover_in(self, e=None):
        if self.disabled_reason:
            self._show_tooltip()
        elif not self.selected:
            self.configure(fg_color="#252530", border_color=ACCENT2)

    def _hover_out(self, e=None):
        self._hide_tooltip()
        if not self.selected and not self.disabled_reason:
            self.configure(fg_color=CARD_BG, border_color="#2e2e3a")

    def _show_tooltip(self):
        if self._tooltip_win:
            return
        x = self.winfo_rootx()
        y = self.winfo_rooty() - 70
        self._tooltip_win = ctk.CTkToplevel(self)
        self._tooltip_win.wm_overrideredirect(True)
        self._tooltip_win.wm_geometry(f"+{x}+{y}")
        self._tooltip_win.configure(fg_color="#1a1a24")
        self._tooltip_win.attributes("-topmost", True)

        # Icon row
        ctk.CTkLabel(
            self._tooltip_win,
            text=self.disabled_reason,
            font=("Segoe UI", 11),
            text_color=WARNING,
            justify="left",
            padx=12, pady=8,
        ).pack()

        # Install hint
        tool = "libreoffice.org" if "LibreOffice" in self.disabled_reason else "calibre-ebook.com"
        ctk.CTkLabel(
            self._tooltip_win,
            text=f"Free download: {tool}",
            font=("Segoe UI", 10),
            text_color=TEXT_DIM,
            padx=12, pady=(0, 8),
        ).pack()

    def _hide_tooltip(self):
        if self._tooltip_win:
            self._tooltip_win.destroy()
            self._tooltip_win = None

    def set_selected(self, selected):
        if self.disabled_reason:
            return
        self.selected = selected
        if selected:
            self.configure(fg_color=ACCENT, border_color=ACCENT)
            self.lbl.configure(text_color="#ffffff", font=("Segoe UI", 13, "bold"))
            self.sub_lbl.configure(text_color="#ddd8ff", font=("Consolas", 12, "bold"))
        else:
            self.configure(fg_color=CARD_BG, border_color="#2e2e3a")
            self.lbl.configure(text_color="#d0d0e8", font=("Segoe UI", 13, "bold"))
            self.sub_lbl.configure(text_color="#9090b8", font=("Consolas", 12, "bold"))


class AnimatedProgressBar(ctk.CTkFrame):
    def __init__(self, master, **kwargs):
        super().__init__(master, fg_color=CARD_BG, corner_radius=6, height=6, **kwargs)
        self.bar        = ctk.CTkFrame(self, fg_color=ACCENT, corner_radius=6, height=6)
        self._animating = False
        self._anim_pos  = 0

    def start_pulse(self):
        self._animating = True
        self._anim_pos  = 0
        self._pulse()

    def _pulse(self):
        if not self._animating:
            return
        w = 0.35
        x = max(0, min((self._anim_pos % 1.6) - 0.4, 1 - w))
        self.bar.place(relx=x, rely=0, relwidth=w, relheight=1)
        self._anim_pos += 0.025
        self.after(16, self._pulse)

    def stop_pulse(self):
        self._animating = False

    def flash_success(self):
        self.bar.configure(fg_color=SUCCESS)
        self.bar.place(relx=0, rely=0, relwidth=1, relheight=1)
        self.after(2000, lambda: self.bar.configure(fg_color=ACCENT))

    def flash_error(self):
        self.bar.configure(fg_color=ERROR)
        self.bar.place(relx=0, rely=0, relwidth=1, relheight=1)
        self.after(2000, lambda: self.bar.configure(fg_color=ACCENT))


# ---------------------------------------------------------------------------
# Main application
# ---------------------------------------------------------------------------

class RecastApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("Recast")
        self.geometry("900x620")
        self.minsize(860, 560)
        self.configure(fg_color=BG)

        self.file_path    = None
        self.category     = None
        self.selected_fmt = None
        self.format_chips = []
        self.active_cat   = "all"

        self._build_layout()
        self._build_sidebar()
        self._build_main()

    def _build_layout(self):
        self.sidebar = ctk.CTkFrame(self, fg_color=SIDEBAR_BG, width=200, corner_radius=0)
        self.sidebar.pack(side="left", fill="y")
        self.sidebar.pack_propagate(False)

        self.main = ctk.CTkFrame(self, fg_color=BG, corner_radius=0)
        self.main.pack(side="left", fill="both", expand=True)

    def _build_sidebar(self):
        logo_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent", height=70)
        logo_frame.pack(fill="x")
        logo_frame.pack_propagate(False)

        ctk.CTkLabel(logo_frame, text="⬡", font=("Segoe UI Emoji", 26),
                     text_color=ACCENT).place(relx=0.18, rely=0.5, anchor="center")
        ctk.CTkLabel(logo_frame, text="Recast", font=("Segoe UI", 15, "bold"),
                     text_color=TEXT_PRI).place(relx=0.58, rely=0.5, anchor="center")

        ctk.CTkFrame(self.sidebar, fg_color=BORDER, height=1).pack(fill="x", padx=16, pady=(0, 10))

        ctk.CTkLabel(self.sidebar, text="CATEGORIES",
                     font=("Segoe UI", 9, "bold"), text_color=TEXT_DIM).pack(
            anchor="w", padx=20, pady=(6, 4)
        )

        self.sidebar_btns = {}
        for cat in CATEGORIES:
            btn = SidebarButton(self.sidebar, cat["icon"], cat["label"],
                                cat["id"], self._on_category_click)
            btn.pack(fill="x")
            self.sidebar_btns[cat["id"]] = btn

        self.sidebar_btns["all"].set_active(True)

        ctk.CTkLabel(self.sidebar, text="Recast v1.0  ·  100% local",
                     font=("Segoe UI", 9), text_color=TEXT_DIM).pack(side="bottom", pady=16)

    def _on_category_click(self, cat_id):
        for k, btn in self.sidebar_btns.items():
            btn.set_active(k == cat_id)
        self.active_cat = cat_id

        # If user clicks a different category while a file is loaded, reset
        if self.file_path and self.category and cat_id not in ("all", self.category):
            self._reset_file()

    def _reset_file(self):
        """Clear the current file and return the UI to its initial state."""
        self.file_path    = None
        self.category     = None
        self.selected_fmt = None
        self.format_chips = []

        # Reset drop zone
        self.drop_frame.configure(border_color=BORDER)
        self.drop_icon.configure(text="⬆", text_color=TEXT_DIM)
        self.drop_title.configure(text="Click to select a file", text_color=TEXT_SEC)
        self.drop_sub.configure(
            text="Images · Documents · Audio · Video · Ebooks · Archives",
            text_color=TEXT_DIM,
        )

        # Hide info strip and format section
        self.info_frame.pack_forget()
        self.format_section.pack_forget()

        # Reset button and status
        self.convert_btn.configure(state="disabled", text="Recast File →")
        self.set_status("Select a file to recast", TEXT_DIM)
        self.prog.stop_pulse()

    def _build_main(self):
        header = ctk.CTkFrame(self.main, fg_color="transparent", height=64)
        header.pack(fill="x", padx=30, pady=(20, 0))
        header.pack_propagate(False)

        ctk.CTkLabel(header, text="Recast", font=FONT_HEAD, text_color=TEXT_PRI).pack(side="left")

        # Drop zone
        self.drop_frame = ctk.CTkFrame(
            self.main, fg_color=CARD_BG, corner_radius=16,
            border_width=2, border_color=BORDER, height=160,
        )
        self.drop_frame.pack(fill="x", padx=30, pady=(16, 0))
        self.drop_frame.pack_propagate(False)

        self.drop_inner = ctk.CTkFrame(self.drop_frame, fg_color="transparent")
        self.drop_inner.place(relx=0.5, rely=0.5, anchor="center")

        self.drop_icon  = ctk.CTkLabel(self.drop_inner, text="⬆",
                                        font=("Segoe UI", 32), text_color=TEXT_DIM)
        self.drop_icon.pack()

        self.drop_title = ctk.CTkLabel(self.drop_inner, text="Click to select a file",
                                        font=FONT_TITLE, text_color=TEXT_SEC)
        self.drop_title.pack()

        self.drop_sub   = ctk.CTkLabel(
            self.drop_inner,
            text="Images · Documents · Audio · Video · Ebooks · Archives",
            font=FONT_SMALL, text_color=TEXT_DIM,
        )
        self.drop_sub.pack()

        for w in [self.drop_frame, self.drop_inner, self.drop_icon, self.drop_title, self.drop_sub]:
            w.configure(cursor="hand2")
            w.bind("<Button-1>", lambda e: self.browse_file())
            w.bind("<Enter>",    self._drop_hover_in)
            w.bind("<Leave>",    self._drop_hover_out)

        # File info strip (shown after selection)
        self.info_frame = ctk.CTkFrame(self.main, fg_color=CARD_BG, corner_radius=12, height=50)

        self.info_name = ctk.CTkLabel(self.info_frame, text="", font=FONT_BODY, text_color=TEXT_PRI)
        self.info_name.place(relx=0.02, rely=0.5, anchor="w")

        self.info_meta = ctk.CTkLabel(self.info_frame, text="", font=FONT_SMALL, text_color=TEXT_DIM)
        self.info_meta.place(relx=0.98, rely=0.5, anchor="e")

        self.change_btn = ctk.CTkButton(
            self.info_frame, text="Change file",
            width=90, height=28, corner_radius=6,
            fg_color=BORDER, hover_color=CARD_HOVER,
            text_color=TEXT_SEC, font=FONT_SMALL,
            command=self.browse_file,
        )
        self.change_btn.place(relx=0.75, rely=0.5, anchor="center")

        # Format picker (shown after selection)
        self.format_section = ctk.CTkFrame(self.main, fg_color="transparent")

        ctk.CTkLabel(self.format_section, text="Output Format",
                     font=("Segoe UI", 14, "bold"), text_color="#c0c0dc").pack(anchor="w", pady=(0, 10))

        self.chips_scroll = ctk.CTkScrollableFrame(
            self.format_section, fg_color="transparent",
            height=115, orientation="horizontal",
        )
        self.chips_scroll.pack(fill="x")

        # Bottom bar
        bottom = ctk.CTkFrame(self.main, fg_color="transparent")
        bottom.pack(side="bottom", fill="x", padx=30, pady=20)

        self.status_label = ctk.CTkLabel(
            bottom, text="Select a file to recast",
            font=FONT_SMALL, text_color=TEXT_DIM,
        )
        self.status_label.pack(anchor="w", pady=(0, 8))

        self.prog = AnimatedProgressBar(bottom)
        self.prog.pack(fill="x", pady=(0, 14))

        self.convert_btn = ctk.CTkButton(
            bottom, text="Recast File →",
            height=50, corner_radius=12,
            fg_color=ACCENT, hover_color="#5a52d5",
            text_color="#ffffff",
            font=("Segoe UI", 15, "bold"),
            command=self.start_conversion,
            state="disabled",
        )
        self.convert_btn.pack(fill="x")

    def _drop_hover_in(self, e=None):
        if not self.file_path:
            self.drop_frame.configure(border_color=ACCENT)

    def _drop_hover_out(self, e=None):
        if not self.file_path:
            self.drop_frame.configure(border_color=BORDER)

    def browse_file(self):
        path = filedialog.askopenfilename()
        if not path:
            return

        self.file_path = path
        ext            = os.path.splitext(path)[1]
        self.category  = get_category(ext)

        if not self.category:
            self.set_status(f"Unsupported file type: {ext}", ERROR)
            return

        self.drop_frame.configure(border_color=SUCCESS)
        self.drop_icon.configure(text="✓", text_color=SUCCESS)
        self.drop_title.configure(text="File ready", text_color=SUCCESS)
        self.drop_sub.configure(text="Click 'Change file' to pick another", text_color=TEXT_DIM)

        filename = os.path.basename(path)
        size     = format_filesize(os.path.getsize(path))
        short    = filename if len(filename) < 42 else filename[:39] + "…"
        self.info_name.configure(text=f"  {short}")
        self.info_meta.configure(text=f"{ext.upper().lstrip('.')}  ·  {size}  ")
        self.info_frame.pack(fill="x", padx=30, pady=(12, 0))

        self.format_section.pack(fill="x", padx=30, pady=(18, 0))
        self._build_format_chips()
        self._on_category_click(self.category)
        self.set_status(f"Detected: {self.category.capitalize()}  ·  {filename}", TEXT_SEC)

    def _build_format_chips(self):
        for w in self.chips_scroll.winfo_children():
            w.destroy()
        self.format_chips.clear()
        self.selected_fmt = None
        self.convert_btn.configure(state="disabled")

        if not self.category:
            return

        in_ext = os.path.splitext(self.file_path)[1].lower().lstrip(".")
        groups = FORMAT_GROUPS.get(self.category, {})
        col    = 0

        for group_label, formats in groups.items():
            valid = [f for f in formats if f in OUTPUT_FORMATS[self.category] and f != in_ext]
            if not valid:
                continue

            ctk.CTkLabel(self.chips_scroll, text=group_label,
                         font=("Segoe UI", 9, "bold"), text_color="#7878a0").grid(
                row=0, column=col, padx=(8, 4), pady=(0, 4), sticky="s"
            )
            col += 1

            for fmt in valid:
                reason = get_chip_disabled_reason(self.category, in_ext, fmt)
                chip = FormatChip(self.chips_scroll, fmt, self._select_format,
                                  disabled_reason=reason)
                chip.grid(row=0, column=col, padx=3, pady=4, sticky="n")
                self.format_chips.append(chip)
                col += 1

            ctk.CTkFrame(self.chips_scroll, fg_color=BORDER, width=1).grid(
                row=0, column=col, padx=8, pady=8, sticky="ns"
            )
            col += 1

        # Auto-select first enabled chip
        first_enabled = next((c for c in self.format_chips if not c.disabled_reason), None)
        if first_enabled:
            self._select_format(first_enabled.fmt)

    def _select_format(self, fmt):
        self.selected_fmt = fmt
        for chip in self.format_chips:
            chip.set_selected(chip.fmt == fmt)
        self.convert_btn.configure(state="normal", text=f"Recast to .{fmt.upper()} →")

    def start_conversion(self):
        if not self.file_path or not self.selected_fmt:
            return

        out_fmt  = self.selected_fmt
        base     = os.path.splitext(self.file_path)[0]
        dst_path = filedialog.asksaveasfilename(
            defaultextension=f".{out_fmt}",
            initialfile=f"{os.path.basename(base)}_recast.{out_fmt}",
            filetypes=[(out_fmt.upper(), f"*.{out_fmt}")],
        )
        if not dst_path:
            return

        self.convert_btn.configure(state="disabled", text="Recasting…")
        self.prog.start_pulse()
        self.set_status("Recasting…  Please wait.", TEXT_SEC)

        threading.Thread(target=self._run, args=(dst_path, out_fmt), daemon=True).start()

    def _run(self, dst_path, out_fmt):
        try:
            in_fmt = os.path.splitext(self.file_path)[1].lower().lstrip(".")
            if self.category == "image":
                convert_image(self.file_path, dst_path, out_fmt)
            elif self.category == "audio":
                convert_audio(self.file_path, dst_path, out_fmt)
            elif self.category == "video":
                convert_video(self.file_path, dst_path, out_fmt)
            elif self.category == "document":
                convert_document(self.file_path, dst_path, in_fmt, out_fmt)
            elif self.category == "ebook":
                convert_ebook(self.file_path, dst_path)
            elif self.category == "archive":
                convert_archive(self.file_path, dst_path, in_fmt, out_fmt)
            self.after(0, self._on_success, dst_path, out_fmt)
        except Exception as e:
            self.after(0, self._on_error, str(e))

    def _on_success(self, dst_path, out_fmt):
        self.prog.stop_pulse()
        self.prog.flash_success()
        self.convert_btn.configure(state="normal", text=f"Recast to .{out_fmt.upper()} →")
        size = format_filesize(os.path.getsize(dst_path))
        self.set_status(f"✓  Saved  ·  {os.path.basename(dst_path)}  ·  {size}", SUCCESS)

    def _on_error(self, msg):
        self.prog.stop_pulse()
        self.prog.flash_error()
        self.convert_btn.configure(state="normal",
                                    text=f"Recast to .{self.selected_fmt.upper()} →")
        self.set_status(f"Error: {msg[:80]}", ERROR)
        messagebox.showerror("Conversion Error", msg)

    def set_status(self, msg, color=None):
        self.status_label.configure(text=msg, text_color=color or TEXT_DIM)


if __name__ == "__main__":
    app = RecastApp()
    app.mainloop()
